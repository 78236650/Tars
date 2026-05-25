"""M3/M4 知识库测试：结构解析、分块、browse 搜索、reindex。"""
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).parent.parent))

FIXTURES = Path(__file__).parent / "fixtures" / "knowledge"


@pytest.fixture
def policy_docx(tmp_path):
    from docx import Document

    path = tmp_path / "policy.docx"
    doc = Document()
    doc.add_heading("总则", level=1)
    doc.add_paragraph("本制度适用于全体仓储管理人员。")
    doc.add_heading("入库管理", level=1)
    doc.add_paragraph("入库须填写申请表并由主管审批。")
    doc.add_heading("出库管理", level=1)
    doc.add_paragraph("出库须双人复核。")
    doc.save(str(path))
    return str(path)


@pytest.fixture
def metrics_xlsx(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")

    path = tmp_path / "metrics.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "KPI"
    ws.append(["指标", "口径"])
    ws.append(["GMV", "含税销售额"])
    ws2 = wb.create_sheet("Cost")
    ws2.append(["项目", "金额"])
    ws2.append(["物流", 100])
    wb.save(str(path))
    return str(path)


@pytest.fixture
def proposal_pptx(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    path = tmp_path / "proposal.pptx"
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "项目背景"
    slide1.placeholders[1].text = "市场需求增长迅速。"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "方案要点"
    slide2.placeholders[1].text = "分三阶段实施。"
    prs.save(str(path))
    return str(path)


def test_parse_policy_docx_sections(policy_docx):
    from tars.knowledge.structure_parser import parse_to_document

    parsed = parse_to_document(policy_docx, doc_type="policy")
    assert len(parsed.sections) >= 3
    titles = [s.title for s in parsed.sections]
    assert "总则" in titles
    assert "入库管理" in titles


def test_parse_metrics_xlsx_sheets(metrics_xlsx):
    from tars.knowledge.structure_parser import parse_to_document

    parsed = parse_to_document(metrics_xlsx)
    assert len(parsed.metrics_tables) == 2
    assert len(parsed.sections) == 2
    assert parsed.doc_type_hint == "metrics"


def test_parse_proposal_pptx_slides(proposal_pptx):
    from tars.knowledge.structure_parser import parse_to_document

    parsed = parse_to_document(proposal_pptx)
    assert len(parsed.sections) >= 2
    assert any("背景" in s.title for s in parsed.sections)


def test_parse_generic_md_headings():
    from tars.knowledge.structure_parser import parse_to_document

    parsed = parse_to_document(str(FIXTURES / "generic.md"))
    assert len(parsed.sections) >= 2
    assert any("概述" in s.title for s in parsed.sections)


def test_chunk_by_sections_respects_boundaries(policy_docx):
    from tars.knowledge.chunker import DocumentChunker
    from tars.knowledge.config import get_chunk_profile
    from tars.knowledge.structure_parser import parse_to_document

    parsed = parse_to_document(policy_docx, doc_type="policy")
    profile = get_chunk_profile("policy")
    chunker = DocumentChunker(chunk_size=profile["chunk_size"], chunk_overlap=profile["overlap"])
    chunks = chunker.chunk_by_sections(parsed.sections)
    assert chunks
    for piece in chunks:
        meta = piece["metadata"]
        assert meta.get("section_id")
        sid = meta["section_id"]
        section = next(s for s in parsed.sections if s.section_id == sid)
        assert piece["text"] in section.text or section.text.startswith(piece["text"][: min(20, len(piece["text"]))])


def test_browse_mode_prefers_summary_chunks():
    from tars.knowledge.access import BROWSE_CHUNK_TYPES, _apply_browse_filter, merge_results_by_doc

    results = [
        {"id": "d1_chunk_0", "score": 0.9, "text": "passage", "metadata": {"chunk_type": "passage", "doc_id": "d1"}},
        {"id": "d1_summary_0", "score": 0.7, "text": "summary", "metadata": {"chunk_type": "doc_summary", "doc_id": "d1"}},
        {"id": "d2_summary_0", "score": 0.8, "text": "s2", "metadata": {"chunk_type": "doc_summary", "doc_id": "d2"}},
    ]
    filtered = _apply_browse_filter(results)
    assert all(r["metadata"]["chunk_type"] in BROWSE_CHUNK_TYPES for r in filtered)
    merged = merge_results_by_doc(filtered, top_k=2)
    assert len(merged) == 2
    assert merged[0]["metadata"]["doc_id"] == "d2"


def test_reindex_estimate_api(tmp_path):
    from fastapi import FastAPI
    from fastapi.testclient import TestClient

    from tars.api.knowledge import init_knowledge_api, router as knowledge_router
    from tars.database import Database

    db = Database(str(tmp_path / "kb.db"))
    conn = db._get_conn()
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO document_collections (id, tenant_id, name, description, created_at, updated_at) "
        "VALUES (?,?,?,?,?,?)",
        ("c1", "default", "t", "", "2026-05-24", "2026-05-24"),
    )
    for i in range(3):
        cur.execute(
            "INSERT INTO document_files (id, collection_id, file_name, file_path, file_type, status, created_at) "
            "VALUES (?,?,?,?,?,?,?)",
            (f"d{i}", "c1", f"f{i}.txt", f"/tmp/f{i}.txt", ".txt", "ready", "2026-05-24"),
        )
    conn.commit()

    app = FastAPI()
    app.include_router(knowledge_router)
    init_knowledge_api(db, vector_store=None, embedding_provider=None)

    with TestClient(app) as client:
        resp = client.post("/api/knowledge/collections/c1/reindex/estimate", json={})
        assert resp.status_code == 200
        data = resp.json()
        assert data["doc_count"] == 3
        assert data["est_tokens"] > 0
