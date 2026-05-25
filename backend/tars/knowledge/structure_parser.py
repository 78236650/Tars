"""结构感知解析入口 — v4.4 知识库深度入库。"""
from __future__ import annotations

import os
import re
from typing import List, Optional

from .metrics_extractor import extract_tables_from_file
from .models import DocumentSection, MetricsTable, ParsedDocument
from .parsers import DocumentParser, _read_text_file

_POLICY_HINTS = ("制度", "规范", "办法", "policy", "regulation", "管理规定")
_PROPOSAL_HINTS = ("方案", "汇报", "提案", "proposal", "deck", "pitch")
_METRICS_HINTS = ("指标", "口径", "metric", "kpi")

_default_parser = DocumentParser()
_HEADING_STYLE_RE = re.compile(r"heading\s*(\d+)", re.I)
_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")


def infer_doc_type(file_name: str, ext: str, explicit: Optional[str] = None) -> str:
    if explicit and explicit.strip().lower() in ("policy", "proposal", "metrics", "generic"):
        return explicit.strip().lower()
    lower_ext = (ext or "").lower()
    if lower_ext in (".xlsx", ".csv"):
        return "metrics"
    if lower_ext == ".pptx":
        return "proposal"
    lname = (file_name or "").lower()
    if any(h in lname for h in _METRICS_HINTS):
        return "metrics"
    if any(h in lname for h in _POLICY_HINTS):
        return "policy"
    if any(h in lname for h in _PROPOSAL_HINTS):
        return "proposal"
    return "generic"


def _heading_level(style_name: str) -> Optional[int]:
    if not style_name:
        return None
    m = _HEADING_STYLE_RE.search(style_name)
    if m:
        return int(m.group(1))
    if style_name.strip().lower() in ("title", "subtitle"):
        return 1 if style_name.lower() == "title" else 2
    return None


def _parse_docx(file_path: str, resolved_type: str) -> ParsedDocument:
    from docx import Document

    doc = Document(file_path)
    sections: List[DocumentSection] = []
    plain_parts: List[str] = []
    current: Optional[DocumentSection] = None
    sec_idx = 0
    warnings: List[str] = []

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style_name = getattr(getattr(para, "style", None), "name", "") or ""
        level = _heading_level(style_name)
        if level is not None:
            if current and current.text.strip():
                sections.append(current)
            current = DocumentSection(
                section_id=f"s{sec_idx}",
                title=text,
                level=level,
                text="",
            )
            sec_idx += 1
            plain_parts.append(text)
        else:
            if current is None:
                current = DocumentSection(section_id="s0", title=os.path.basename(file_path), level=0, text="")
            current.text = (current.text + "\n" + text).strip()
            plain_parts.append(text)

    if current and current.text.strip():
        sections.append(current)

    if not sections:
        warnings.append("docx_no_headings")
        fallback = "\n".join(p.text.strip() for p in doc.paragraphs if p.text and p.text.strip())
        if fallback:
            sections = [DocumentSection(section_id="s0", title=os.path.basename(file_path), level=0, text=fallback)]
            plain_parts = [fallback]

    return ParsedDocument(
        plain_text="\n".join(plain_parts),
        sections=sections,
        doc_type_hint=resolved_type,
        parse_warnings=warnings,
    )


def _parse_pptx(file_path: str, resolved_type: str) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as e:
        return ParsedDocument(
            plain_text="",
            sections=[],
            doc_type_hint=resolved_type,
            parse_warnings=[f"pptx_import_failed: {e}"],
        )

    prs = Presentation(file_path)
    sections: List[DocumentSection] = []
    plain_parts: List[str] = []
    warnings: List[str] = []

    for i, slide in enumerate(prs.slides):
        texts: List[str] = []
        title = f"Slide {i + 1}"
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if not t:
                continue
            texts.append(t)
            try:
                if getattr(shape, "is_placeholder", False) and shape.placeholder_format.type == 1:
                    title = t
            except Exception:
                pass
        body = "\n".join(texts)
        if not body.strip():
            warnings.append(f"slide_{i + 1}_empty")
            continue
        sections.append(
            DocumentSection(
                section_id=f"slide_{i}",
                title=title,
                level=1,
                text=body,
                page_or_slide=i + 1,
            )
        )
        plain_parts.append(body)

    return ParsedDocument(
        plain_text="\n\n".join(plain_parts),
        sections=sections,
        doc_type_hint=resolved_type,
        parse_warnings=warnings,
    )


def _build_parsed_from_tables(tables: List[MetricsTable], resolved_type: str) -> ParsedDocument:
    sections: List[DocumentSection] = []
    plain_parts: List[str] = []
    for i, table in enumerate(tables):
        lines = []
        if table.headers:
            lines.append(" | ".join(table.headers))
        for row in table.rows[:200]:
            lines.append(" | ".join(str(row.get(h, "")) for h in table.headers))
        body = "\n".join(lines)
        sections.append(
            DocumentSection(
                section_id=f"sheet_{i}",
                title=table.sheet_name or f"Sheet {i + 1}",
                level=1,
                text=body,
            )
        )
        plain_parts.append(body)
    return ParsedDocument(
        plain_text="\n\n".join(plain_parts),
        sections=sections,
        doc_type_hint=resolved_type,
        metrics_tables=tables,
        parse_warnings=[] if tables else ["tables_empty"],
    )


def _parse_xlsx(file_path: str, resolved_type: str) -> ParsedDocument:
    tables = extract_tables_from_file(file_path)
    return _build_parsed_from_tables(tables, resolved_type)


def _parse_csv(file_path: str, resolved_type: str) -> ParsedDocument:
    tables = extract_tables_from_file(file_path)
    if tables:
        return _build_parsed_from_tables(tables, resolved_type)
    return _parse_plain(file_path, resolved_type)


def _parse_markdown(file_path: str, resolved_type: str) -> ParsedDocument:
    raw = _read_text_file(file_path)
    sections: List[DocumentSection] = []
    plain_parts: List[str] = []
    current: Optional[DocumentSection] = None
    sec_idx = 0

    for line in raw.splitlines():
        m = _MD_HEADING_RE.match(line)
        if m:
            if current and current.text.strip():
                sections.append(current)
            level = len(m.group(1))
            title = m.group(2).strip()
            current = DocumentSection(section_id=f"s{sec_idx}", title=title, level=level, text="")
            sec_idx += 1
            plain_parts.append(title)
        else:
            if current is None:
                current = DocumentSection(section_id="s0", title=os.path.basename(file_path), level=0, text="")
            if line.strip():
                current.text = (current.text + "\n" + line).strip()
                plain_parts.append(line.strip())

    if current and (current.text.strip() or current.title):
        if not current.text.strip() and len(sections) > 0:
            pass
        else:
            sections.append(current)

    if not sections:
        text = raw.strip()
        if text:
            sections = [DocumentSection(section_id="s0", title=os.path.basename(file_path), level=0, text=text)]
            plain_parts = [text]

    return ParsedDocument(
        plain_text="\n".join(plain_parts),
        sections=sections,
        doc_type_hint=resolved_type,
        parse_warnings=[],
    )


def _parse_plain(file_path: str, resolved_type: str) -> ParsedDocument:
    text = _default_parser.parse(file_path).strip()
    sections = []
    if text:
        sections = [DocumentSection(section_id="s0", title=os.path.basename(file_path), level=0, text=text)]
    return ParsedDocument(plain_text=text, sections=sections, doc_type_hint=resolved_type, parse_warnings=[])


def parse_to_document(file_path: str, doc_type: Optional[str] = None) -> ParsedDocument:
    """按扩展名路由到结构化解析器，返回 ParsedDocument。"""
    file_name = os.path.basename(file_path)
    ext = os.path.splitext(file_name)[1].lower()
    resolved_type = infer_doc_type(file_name, ext, doc_type)

    try:
        if ext == ".docx":
            return _parse_docx(file_path, resolved_type)
        if ext == ".pptx":
            return _parse_pptx(file_path, resolved_type)
        if ext == ".xlsx":
            return _parse_xlsx(file_path, resolved_type)
        if ext == ".csv":
            parsed = _parse_csv(file_path, resolved_type)
            if parsed.metrics_tables:
                return parsed
        if ext == ".md":
            return _parse_markdown(file_path, resolved_type)
        return _parse_plain(file_path, resolved_type)
    except Exception as e:
        return ParsedDocument(
            plain_text="",
            sections=[],
            doc_type_hint=resolved_type,
            parse_warnings=[f"parse_failed: {e}"],
        )
