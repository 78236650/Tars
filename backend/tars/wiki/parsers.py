"""Wiki 文档解析器 — 轻量版。"""
from __future__ import annotations

import csv, os, re

def _read_text_file(file_path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "gb18030", "gbk"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def parse_pdf(file_path: str) -> str:
    import fitz
    return "\n".join(p.get_text("text") for p in fitz.open(file_path)).strip()

def parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for t in doc.tables:
        for r in t.rows:
            cells = [c.text.strip() for c in r.cells if c.text.strip()]
            if cells: parts.append(" | ".join(cells))
    return "\n".join(parts)

def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    parts = []
    for s in Presentation(file_path).slides:
        for sh in s.shapes:
            if hasattr(sh, "text") and sh.text:
                parts.append(sh.text.strip())
    return "\n".join(parts)

def parse_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells: lines.append(", ".join(cells))
    wb.close()
    return "\n".join(lines)

_PARSERS = {
    ".txt": lambda p: _read_text_file(p),
    ".md": lambda p: _read_text_file(p),
    ".pdf": parse_pdf, ".docx": parse_docx, ".pptx": parse_pptx,
    ".xlsx": parse_xlsx, ".csv": lambda p: _read_text_file(p),
}

def parse_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    parser = _PARSERS.get(ext)
    if not parser:
        raise ValueError(f"不支持的文档类型: {ext or 'unknown'}")
    return parser(file_path)
